"""
Generate a professional 15-slide PowerPoint presentation for NST Studio.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

# ── Color palette ─────────────────────────────────────────────────────────────
BG_DARK    = RGBColor(0x0D, 0x0D, 0x1A)
BG_CARD    = RGBColor(0x16, 0x16, 0x2E)
BG_CODE    = RGBColor(0x0A, 0x0A, 0x18)
ACCENT     = RGBColor(0x7C, 0x3A, 0xED)
ACCENT2    = RGBColor(0xEC, 0x4E, 0x20)
ACCENT3    = RGBColor(0x10, 0x99, 0x73)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xB0, 0xB8, 0xCC)
YELLOW     = RGBColor(0xF5, 0xC5, 0x18)
BLUE_CODE  = RGBColor(0x7D, 0xD3, 0xFC)
GREEN_CODE = RGBColor(0x86, 0xEF, 0xAC)
PINK       = RGBColor(0xEC, 0x48, 0x99)
ORANGE     = RGBColor(0xF5, 0x9E, 0x0B)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

PROJECT_DIR = os.path.dirname(os.path.abspath(__file__))
STYLE_IMG   = os.path.join(PROJECT_DIR, "NST/images/styles/pencil_style_img1.jpg")
CONTENT_IMG = os.path.join(PROJECT_DIR, "NST/images/content/content_img1.png")

prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H
BLANK = prs.slide_layouts[6]

TOTAL = 15


# ── Helpers ───────────────────────────────────────────────────────────────────

def add_bg(slide, color=BG_DARK):
    bg = slide.shapes.add_shape(1, 0, 0, SLIDE_W, SLIDE_H)
    bg.fill.solid()
    bg.fill.fore_color.rgb = color
    bg.line.fill.background()

def add_rect(slide, l, t, w, h, color):
    shape = slide.shapes.add_shape(1, l, t, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_text(slide, text, l, t, w, h,
             size=20, bold=False, color=WHITE,
             align=PP_ALIGN.LEFT, wrap=True, italic=False):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return tb

def add_bullets(slide, items, l, t, w, h,
                size=16, color=LIGHT_GRAY, title=None, title_color=ACCENT,
                title_size=None):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    first = True
    if title:
        p = tf.paragraphs[0]
        first = False
        run = p.add_run()
        run.text = title
        run.font.size = Pt(title_size or size + 2)
        run.font.bold = True
        run.font.color.rgb = title_color
    for item in items:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_before = Pt(5)
        run = p.add_run()
        run.text = f"  •  {item}"
        run.font.size = Pt(size)
        run.font.color.rgb = color

def accent_bar(slide, color=ACCENT):
    bar = slide.shapes.add_shape(1, 0, Inches(0.08), SLIDE_W, Pt(4))
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()

def slide_num(slide, n):
    add_text(slide, f"{n} / {TOTAL}",
             Inches(12.1), Inches(7.1), Inches(1.1), Inches(0.35),
             size=11, color=LIGHT_GRAY, align=PP_ALIGN.RIGHT)

def title_text(slide, text, subtitle=None):
    add_text(slide, text,
             Inches(0.5), Inches(0.22), Inches(12.2), Inches(0.65),
             size=34, bold=True, color=WHITE)
    if subtitle:
        add_text(slide, subtitle,
                 Inches(0.5), Inches(0.88), Inches(12.2), Inches(0.38),
                 size=14, color=LIGHT_GRAY, italic=True)

def add_image_safe(slide, path, l, t, w, h):
    if os.path.exists(path):
        slide.shapes.add_picture(path, l, t, w, h)
    else:
        add_rect(slide, l, t, w, h, RGBColor(0x25, 0x25, 0x40))
        add_text(slide, "[ image ]", l, t + h // 3, w, Inches(0.4),
                 size=12, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

def card(slide, l, t, w, h, left_color=None):
    add_rect(slide, l, t, w, h, BG_CARD)
    if left_color:
        add_rect(slide, l, t, Inches(0.08), h, left_color)


# ══════════════════════════════════════════════════════════════════════════════
# 1 — Title
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_bg(s)
add_rect(s, 0, 0, Inches(0.12), SLIDE_H, ACCENT)

add_text(s, "NST Studio",
         Inches(0.4), Inches(1.3), Inches(8.5), Inches(1.8),
         size=76, bold=True, color=WHITE)
add_text(s,
    "Neural Style Transfer\nFrom Classic Optimization to Real-Time Inference",
    Inches(0.4), Inches(3.2), Inches(8.5), Inches(1.1),
    size=22, color=LIGHT_GRAY)

for i, (lbl, col) in enumerate([
    ("PyTorch 2.x", ACCENT),
    ("Flask 3.x", ACCENT2),
    ("VGG-19", ACCENT3),
    ("PatchGAN", ORANGE),
]):
    add_rect(s, Inches(0.4 + i * 2.1), Inches(4.55), Inches(1.9), Inches(0.44), col)
    add_text(s, lbl,
             Inches(0.4 + i * 2.1), Inches(0.55 + Inches(4.0)), Inches(1.9), Inches(0.44),
             size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

add_image_safe(s, STYLE_IMG, Inches(9.1), Inches(0.7), Inches(3.9), Inches(5.7))
add_text(s, "Deep Learning Project  ·  2026",
         Inches(0.4), Inches(7.0), Inches(6), Inches(0.38),
         size=13, color=LIGHT_GRAY)
slide_num(s, 1)


# ══════════════════════════════════════════════════════════════════════════════
# 2 — What is Neural Style Transfer?
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_bg(s); accent_bar(s)
title_text(s, "What is Neural Style Transfer?")

add_rect(s, Inches(0.5), Inches(1.1), Inches(12.3), Inches(1.3), BG_CARD)
add_text(s,
    "Neural Style Transfer (NST) blends the content of one image with the "
    "artistic style of another using a deep convolutional neural network — "
    "without any human-crafted filters.",
    Inches(0.65), Inches(1.18), Inches(12.0), Inches(1.1),
    size=18, color=WHITE, wrap=True)

for i, (icon, title, desc) in enumerate([
    ("📷", "Content Image",  "The photograph to preserve — objects, structure, composition"),
    ("🎨", "Style Image",    "The artwork to borrow from — textures, colors, brushstrokes"),
    ("✨", "Stylized Output", "A new image painted in the style while keeping the content"),
]):
    x = Inches(0.5 + i * 4.3)
    card(s, x, Inches(2.6), Inches(4.05), Inches(2.5), ACCENT)
    add_text(s, icon, x + Inches(0.2), Inches(2.72), Inches(0.6), Inches(0.55), size=28)
    add_text(s, title, x + Inches(0.2), Inches(3.3), Inches(3.7), Inches(0.4),
             size=17, bold=True, color=ACCENT)
    add_text(s, desc, x + Inches(0.2), Inches(3.75), Inches(3.7), Inches(0.95),
             size=14, color=LIGHT_GRAY, wrap=True)

add_text(s, "Result  =  f ( Content Image,  Style Image )",
         Inches(0.5), Inches(5.35), Inches(12.3), Inches(0.55),
         size=26, bold=True, color=YELLOW, align=PP_ALIGN.CENTER)

add_bullets(s, [
    "Originally proposed by Gatys, Ecker & Bethge (2016)",
    "No training required for classic NST — only an iterative optimization",
    "Later extended to real-time feed-forward networks by Johnson et al. (2016)",
], Inches(0.5), Inches(6.05), Inches(12.3), Inches(1.2), size=14)
slide_num(s, 2)


# ══════════════════════════════════════════════════════════════════════════════
# 3 — The Problem We're Solving
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_bg(s); accent_bar(s)
title_text(s, "The Problem We're Solving",
           "Classic NST produces great results but is painfully slow")

problems = [
    ("⏱️", "Speed",      "Classic NST takes 5–15 minutes per image\n300–1000 gradient steps through a 143M-param VGG-19"),
    ("🖥️", "Usability",  "Command-line tools are inaccessible to non-technical users\nNo visual feedback during generation"),
    ("🎨", "Flexibility", "Fixed style — changing it means rerunning the whole optimization\nNo way to compare multiple styles instantly"),
    ("📦", "Scale",       "Not practical for batch processing or real-time applications\nEvery new image needs a fresh 10-minute run"),
]
for i, (icon, title, desc) in enumerate(problems):
    col = 0 if i % 2 == 0 else 1
    row = i // 2
    x = Inches(0.5 + col * 6.45)
    y = Inches(1.15 + row * 2.0)
    card(s, x, y, Inches(6.1), Inches(1.8), ACCENT2)
    add_text(s, icon, x + Inches(0.22), y + Inches(0.15), Inches(0.6), Inches(0.6), size=28)
    add_text(s, title, x + Inches(0.9), y + Inches(0.15), Inches(4.9), Inches(0.42),
             size=18, bold=True, color=WHITE)
    add_text(s, desc, x + Inches(0.22), y + Inches(0.72), Inches(5.65), Inches(0.9),
             size=13, color=LIGHT_GRAY, wrap=True)

add_rect(s, Inches(0.5), Inches(5.35), Inches(12.3), Inches(0.72), RGBColor(0x1A, 0x2A, 0x1A))
add_text(s,
    "✅  Our solution:  a full-stack web app with Classic NST for quality  +  "
    "Fast NST for instant real-time inference",
    Inches(0.65), Inches(5.43), Inches(12.0), Inches(0.55),
    size=16, bold=True, color=GREEN_CODE)
slide_num(s, 3)


# ══════════════════════════════════════════════════════════════════════════════
# 4 — How Classic NST Works
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_bg(s); accent_bar(s)
title_text(s, "How Classic NST Works")

steps = [
    ("1", "Freeze VGG-19",       "Load ImageNet-pretrained VGG-19 — weights never update"),
    ("2", "Extract Features",    "Pass content & style images through VGG-19; save layer activations"),
    ("3", "Compute Gram Matrix", "For each style layer: Gram(F) = F·Fᵀ ÷ (C·H·W) — captures texture statistics"),
    ("4", "Initialize Output",   "Start from the content image (or random noise)"),
    ("5", "Optimize Pixels",     "Minimize L_total = α·L_content + β·L_style using L-BFGS / Adam"),
    ("6", "Iterate",             "300–1000 steps — the image pixels are the only learned parameters"),
]
for i, (num, title, desc) in enumerate(steps):
    y = Inches(1.12 + i * 1.0)
    add_rect(s, Inches(0.5), y + Inches(0.07), Inches(0.52), Inches(0.52), ACCENT)
    add_text(s, num, Inches(0.5), y + Inches(0.05), Inches(0.52), Inches(0.52),
             size=17, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    card(s, Inches(1.18), y, Inches(8.8), Inches(0.88))
    add_text(s, title, Inches(1.32), y + Inches(0.06), Inches(3.2), Inches(0.36),
             size=16, bold=True, color=WHITE)
    add_text(s, desc, Inches(1.32), y + Inches(0.44), Inches(8.4), Inches(0.38),
             size=13, color=LIGHT_GRAY, wrap=True)

# VGG layer panel
card(s, Inches(10.25), Inches(1.12), Inches(2.6), Inches(6.0))
add_text(s, "VGG-19 Layers", Inches(10.38), Inches(1.22), Inches(2.35), Inches(0.38),
         size=14, bold=True, color=ACCENT)
for j, (lyr, col) in enumerate([
    ("relu1_2  ← style", ACCENT),
    ("relu2_2  ← style", ACCENT),
    ("relu3_3  ← style", ACCENT),
    ("relu4_2  ← content", YELLOW),
    ("relu4_3  ← style", ACCENT),
    ("relu5_3  ← style", ACCENT),
]):
    add_text(s, lyr, Inches(10.38), Inches(1.72 + j * 0.85), Inches(2.35), Inches(0.4),
             size=12, color=col)
slide_num(s, 4)


# ══════════════════════════════════════════════════════════════════════════════
# 5 — The Math: Loss Functions
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_bg(s); accent_bar(s)
title_text(s, "The Math: Loss Functions")

# Content loss
card(s, Inches(0.5), Inches(1.1), Inches(5.9), Inches(2.7), ACCENT3)
add_text(s, "Content Loss", Inches(0.65), Inches(1.18), Inches(5.6), Inches(0.42),
         size=17, bold=True, color=ACCENT3)
add_text(s,
    "L_content = ½ · Σ (F_ij − P_ij)²\n\n"
    "F = features of generated image\n"
    "P = features of content image\n"
    "Layer: relu4_2 (conv4_2)",
    Inches(0.65), Inches(1.65), Inches(5.6), Inches(1.9),
    size=14, color=LIGHT_GRAY, wrap=True)

# Style loss
card(s, Inches(6.7), Inches(1.1), Inches(6.1), Inches(2.7), ACCENT)
add_text(s, "Style Loss", Inches(6.85), Inches(1.18), Inches(5.8), Inches(0.42),
         size=17, bold=True, color=ACCENT)
add_text(s,
    "L_style = Σ_l  w_l · MSE( G_l,  A_l )\n\n"
    "G_l = Gram matrix of generated image at layer l\n"
    "A_l = Gram matrix of style image at layer l\n"
    "Gram(F) = F·Fᵀ  ÷  (C·H·W)",
    Inches(6.85), Inches(1.65), Inches(5.8), Inches(1.9),
    size=14, color=LIGHT_GRAY, wrap=True)

# Total loss
add_rect(s, Inches(0.5), Inches(4.0), Inches(12.3), Inches(1.1), RGBColor(0x1A, 0x14, 0x2E))
add_text(s, "L_total  =  α · L_content  +  β · L_style",
         Inches(0.5), Inches(4.08), Inches(12.3), Inches(0.55),
         size=28, bold=True, color=YELLOW, align=PP_ALIGN.CENTER)
add_text(s, "α = content weight  (e.g. 1e5)        β = style weight  (e.g. 1e10)",
         Inches(0.5), Inches(4.63), Inches(12.3), Inches(0.38),
         size=14, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

# Gram insight
card(s, Inches(0.5), Inches(5.3), Inches(12.3), Inches(1.9))
add_text(s, "Why Gram Matrices?", Inches(0.65), Inches(5.38), Inches(4), Inches(0.4),
         size=16, bold=True, color=ACCENT)
add_text(s,
    "The Gram matrix captures correlations between feature channels — not where features appear, "
    "but how they co-occur. This makes it position-invariant, exactly what we need for texture "
    "and style: the painting's brushstroke pattern is the same wherever you look.",
    Inches(0.65), Inches(5.82), Inches(12.0), Inches(1.1),
    size=14, color=LIGHT_GRAY, wrap=True)
slide_num(s, 5)


# ══════════════════════════════════════════════════════════════════════════════
# 6 — VGG-19 Feature Extraction
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_bg(s); accent_bar(s)
title_text(s, "VGG-19: Why It Works for Style Transfer")

add_bullets(s, [
    "19-layer deep CNN trained on 1.2M ImageNet images",
    "Learns a rich hierarchy: edges → textures → objects → semantics",
    "Early layers (relu1–2): low-level textures, colors, edges",
    "Middle layers (relu3–4): patterns, brushstrokes, material",
    "Deep layers (relu5): high-level objects and scene structure",
], Inches(0.5), Inches(1.12), Inches(6.3), Inches(3.8),
   size=16, title="Why VGG-19?", title_size=18)

# Architecture strip
layers = [
    ("Conv\n1-2",   "64ch",  ACCENT),
    ("Conv\n3-4",   "128ch", ACCENT),
    ("Conv\n5-7",   "256ch", ACCENT),
    ("Conv\n8-11",  "512ch", ACCENT2),
    ("Conv\n12-16", "512ch", ACCENT2),
    ("FC\n1-3",     "class", LIGHT_GRAY),
]
for i, (lbl, ch, col) in enumerate(layers):
    x = Inches(0.5 + i * 2.1)
    h = Inches(0.8 + i * 0.28)
    y = Inches(5.25) - h
    add_rect(s, x, Inches(5.25) - h, Inches(1.8), h, col)
    add_text(s, lbl, x, Inches(5.25) - h + Inches(0.08), Inches(1.8), Inches(0.55),
             size=11, bold=True, color=BG_DARK, align=PP_ALIGN.CENTER)
    add_text(s, ch, x, Inches(5.25), Inches(1.8), Inches(0.3),
             size=11, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

add_text(s, "← low-level (style)                                  high-level (content) →",
         Inches(0.5), Inches(5.62), Inches(12.3), Inches(0.35),
         size=12, color=LIGHT_GRAY, italic=True, align=PP_ALIGN.CENTER)

# Right side — implementation note
card(s, Inches(7.0), Inches(1.12), Inches(5.8), Inches(4.8), ACCENT)
add_text(s, "Implementation Notes", Inches(7.15), Inches(1.2), Inches(5.5), Inches(0.42),
         size=17, bold=True, color=WHITE)
add_bullets(s, [
    "Weights loaded from torchvision (ImageNet pretrained)",
    "All parameters frozen — requires_grad = False",
    "Custom normalization layer prepended (mean/std of ImageNet)",
    "Content hooks at relu4_2",
    "Style hooks at relu1_2, relu2_2, relu3_3, relu4_3, relu5_3",
    "CRITICAL: must run in float32 — relu4+ activations exceed float16 max (65504)",
], Inches(7.15), Inches(1.7), Inches(5.5), Inches(3.8), size=14)
slide_num(s, 6)


# ══════════════════════════════════════════════════════════════════════════════
# 7 — System Architecture
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_bg(s); accent_bar(s)
title_text(s, "System Architecture")

# Browser
add_rect(s, Inches(0.5), Inches(1.12), Inches(12.3), Inches(1.1), RGBColor(0x1A,0x2A,0x1A))
add_text(s, "🌐  Browser  (HTML + CSS + JavaScript)",
         Inches(0.65), Inches(1.2), Inches(7), Inches(0.38),
         size=16, bold=True, color=GREEN_CODE)
add_text(s, "Drag-drop upload  ·  Before/After slider  ·  Gallery  ·  Fullscreen viewer  ·  Toast notifications",
         Inches(0.65), Inches(1.6), Inches(11.5), Inches(0.38),
         size=13, color=LIGHT_GRAY)

add_text(s, "HTTP / SSE  ▼", Inches(5.7), Inches(2.38), Inches(2.5), Inches(0.38),
         size=13, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

# Flask
add_rect(s, Inches(2.0), Inches(2.82), Inches(9.3), Inches(1.0), RGBColor(0x1A,0x1A,0x2E))
add_text(s, "⚙️  Flask  app.py",
         Inches(2.15), Inches(2.9), Inches(4), Inches(0.38),
         size=16, bold=True, color=ACCENT2)
add_text(s, "/upload  ·  /progress/<id>  ·  /fast-stylize  ·  /gallery  ·  /models  ·  /sample-img",
         Inches(2.15), Inches(3.3), Inches(8.9), Inches(0.35),
         size=13, color=LIGHT_GRAY)

# Two arrows
for x in [Inches(3.2), Inches(8.8)]:
    add_text(s, "▼", x, Inches(3.98), Inches(1), Inches(0.38),
             size=20, color=ACCENT, align=PP_ALIGN.CENTER)

# Two engines
card(s, Inches(0.5), Inches(4.42), Inches(5.9), Inches(1.65), ACCENT3)
add_text(s, "🔬  nst_core.py — Classic NST",
         Inches(0.65), Inches(4.5), Inches(5.65), Inches(0.4),
         size=15, bold=True, color=WHITE)
add_bullets(s, ["VGG-19 pixel optimization", "SSE progress streaming", "L-BFGS / Adam optimizer"],
            Inches(0.65), Inches(4.92), Inches(5.65), Inches(1.1), size=13)

card(s, Inches(6.9), Inches(4.42), Inches(5.9), Inches(1.65), ACCENT)
add_text(s, "⚡  fast_nst.py — Fast NST",
         Inches(7.05), Inches(4.5), Inches(5.65), Inches(0.4),
         size=15, bold=True, color=WHITE)
add_bullets(s, ["TransformerNet feed-forward", "Sub-second inference", "Pre-trained .pth models"],
            Inches(7.05), Inches(4.92), Inches(5.65), Inches(1.1), size=13)

# Storage row
for i, (icon, label, col) in enumerate([
    ("💾", "models/  .pth files", BG_CARD),
    ("🖼️", "static/outputs/  results", BG_CARD),
    ("📊", "runs/  TensorBoard logs", BG_CARD),
    ("🎞️", "samples/  training previews", BG_CARD),
]):
    x = Inches(0.5 + i * 3.22)
    card(s, x, Inches(6.3), Inches(3.0), Inches(0.88))
    add_text(s, f"{icon}  {label}", x + Inches(0.18), Inches(6.42), Inches(2.75), Inches(0.38),
             size=12, color=LIGHT_GRAY)
slide_num(s, 7)


# ══════════════════════════════════════════════════════════════════════════════
# 8 — Flask API
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_bg(s); accent_bar(s)
title_text(s, "Flask REST API Endpoints")

endpoints = [
    ("POST", "/upload",              "Upload content + style images, start Classic NST job, return job ID"),
    ("GET",  "/progress/<id>",       "SSE stream — sends loss value and preview image every N iterations"),
    ("POST", "/fast-stylize",        "One-shot inference with a pre-trained Fast NST TransformerNet model"),
    ("GET",  "/models",              "List available Fast NST .pth models with name, path and thumbnail"),
    ("GET",  "/gallery",             "Return last 60 stylized outputs sorted by modification time"),
    ("GET",  "/sample-img/<m>/<f>",  "Serve per-model training preview images saved during training"),
    ("GET",  "/samples/<model>",     "List all training preview filenames for a given model"),
]
for i, (method, ep, desc) in enumerate(endpoints):
    y = Inches(1.1 + i * 0.86)
    card(s, Inches(0.5), y, Inches(12.3), Inches(0.76))
    add_rect(s, Inches(0.5), y, Inches(0.08), Inches(0.76),
             ACCENT2 if method == "POST" else ACCENT)
    mcol = ACCENT2 if method == "POST" else ACCENT3
    add_text(s, method, Inches(0.72), y + Inches(0.18), Inches(0.82), Inches(0.36),
             size=12, bold=True, color=mcol)
    add_text(s, ep, Inches(1.6), y + Inches(0.18), Inches(3.8), Inches(0.36),
             size=14, bold=True, color=BLUE_CODE)
    add_text(s, desc, Inches(5.6), y + Inches(0.2), Inches(7.1), Inches(0.36),
             size=13, color=LIGHT_GRAY, wrap=True)

add_text(s, "Built with Flask 3.x  ·  Threaded job dict  ·  JSON metadata  ·  Pathlib for all paths",
         Inches(0.5), Inches(7.1), Inches(12.3), Inches(0.3),
         size=12, color=LIGHT_GRAY, italic=True)
slide_num(s, 8)


# ══════════════════════════════════════════════════════════════════════════════
# 9 — Classic NST Core Code
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_bg(s); accent_bar(s)
title_text(s, "Classic NST — Core Implementation  (nst_core.py)")

# Gram matrix code
add_rect(s, Inches(0.5), Inches(1.12), Inches(5.9), Inches(2.5), BG_CODE)
add_text(s, "Gram Matrix", Inches(0.65), Inches(1.18), Inches(5.6), Inches(0.35),
         size=13, bold=True, color=ACCENT)
add_text(s,
    "def gram_matrix(input):\n"
    "    a, b, c, d = input.size()\n"
    "    features = input.view(a * b, c * d)\n"
    "    G = torch.mm(features, features.t())\n"
    "    return G.div(a * b * c * d)",
    Inches(0.65), Inches(1.56), Inches(5.6), Inches(1.75),
    size=13, color=BLUE_CODE, wrap=False)

# Style loss code
add_rect(s, Inches(0.5), Inches(3.75), Inches(5.9), Inches(2.5), BG_CODE)
add_text(s, "Style Loss Module", Inches(0.65), Inches(3.81), Inches(5.6), Inches(0.35),
         size=13, bold=True, color=ACCENT)
add_text(s,
    "class StyleLoss(nn.Module):\n"
    "    def __init__(self, target_feature):\n"
    "        super().__init__()\n"
    "        self.target = gram_matrix(\n"
    "            target_feature).detach()\n"
    "    def forward(self, input):\n"
    "        G = gram_matrix(input)\n"
    "        self.loss = F.mse_loss(G, self.target)\n"
    "        return input",
    Inches(0.65), Inches(4.2), Inches(5.6), Inches(1.85),
    size=12, color=BLUE_CODE, wrap=False)

# Right side notes
add_bullets(s, [
    "ContentLoss: MSE between generated and content features at relu4_2",
    "StyleLoss: MSE between Gram matrices at 5 style layers",
    "Normalization layer prepended so images stay in [0,1]",
    "Loss hooks injected into frozen VGG-19 graph",
    "L-BFGS optimizer acts directly on image tensor",
    "SSE streams loss + base64 preview every 50 steps to browser",
    "Final image decoded from float tensor → PIL → JPEG → disk",
], Inches(6.7), Inches(1.12), Inches(6.1), Inches(5.5),
   size=15, title="Implementation Notes")
slide_num(s, 9)


# ══════════════════════════════════════════════════════════════════════════════
# 10 — Web UI Features
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_bg(s); accent_bar(s)
title_text(s, "Web UI Features")

features = [
    ("🖱️", "Drag & Drop Upload",    "Drop content and style images anywhere — no file picker needed",           ACCENT),
    ("↔️", "Before / After Slider",  "CSS clip-path divider to compare original vs. stylized side-by-side",     ACCENT2),
    ("🖼️", "Results Gallery",        "Last 60 outputs with metadata — model name, style weights, timestamp",    ACCENT3),
    ("⛶",  "Fullscreen Viewer",      "Click any image to expand full-screen, ESC to close",                     ORANGE),
    ("📊", "Live SSE Progress",      "Loss value + preview image streamed live as Classic NST iterates",        PINK),
    ("🔔", "Toast Notifications",    "Non-blocking pop-ups for job complete, errors, and model selection",      YELLOW),
]
for i, (icon, title, desc, col) in enumerate(features):
    c, row = i % 2, i // 2
    x = Inches(0.5 + c * 6.45)
    y = Inches(1.12 + row * 1.85)
    card(s, x, y, Inches(6.1), Inches(1.7), col)
    add_text(s, icon, x + Inches(0.22), y + Inches(0.15), Inches(0.6), Inches(0.6), size=26)
    add_text(s, title, x + Inches(0.9), y + Inches(0.15), Inches(4.95), Inches(0.42),
             size=16, bold=True, color=WHITE)
    add_text(s, desc, x + Inches(0.22), y + Inches(0.75), Inches(5.65), Inches(0.78),
             size=13, color=LIGHT_GRAY, wrap=True)
slide_num(s, 10)


# ══════════════════════════════════════════════════════════════════════════════
# 11 — Dataset & Preprocessing
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_bg(s); accent_bar(s)
title_text(s, "Dataset & Preprocessing")

# Dataset stats
card(s, Inches(0.5), Inches(1.12), Inches(5.9), Inches(3.5), ACCENT3)
add_text(s, "MS-COCO train2014", Inches(0.65), Inches(1.2), Inches(5.65), Inches(0.42),
         size=18, bold=True, color=WHITE)

stats = [
    ("Total images",    "82,783"),
    ("Valid images",    "82,177"),
    ("Removed (corrupt/small)", "606"),
    ("Dataset size",    "~13 GB"),
    ("Min image size",  "256 × 256 px"),
    ("Image format",    "JPEG"),
]
for j, (k, v) in enumerate(stats):
    add_text(s, k, Inches(0.65), Inches(1.72 + j * 0.45), Inches(3.3), Inches(0.38),
             size=14, color=LIGHT_GRAY)
    add_text(s, v, Inches(4.1), Inches(1.72 + j * 0.45), Inches(2.1), Inches(0.38),
             size=14, bold=True, color=YELLOW, align=PP_ALIGN.RIGHT)

# Preprocessing pipeline
card(s, Inches(6.7), Inches(1.12), Inches(6.1), Inches(3.5))
add_text(s, "preprocess.py — Validation Pipeline", Inches(6.85), Inches(1.2), Inches(5.85), Inches(0.42),
         size=16, bold=True, color=ACCENT)
add_bullets(s, [
    "8 parallel workers via ThreadPoolExecutor",
    "Force-decode every pixel — catches truncated JPEGs",
    "Filter images smaller than 256px on either side",
    "Writes valid_images.txt manifest — ~2 min one-time cost",
    "Manifest used by DataLoader, skips re-validation on resume",
], Inches(6.85), Inches(1.68), Inches(5.85), Inches(2.8), size=14)

# Training transform
card(s, Inches(0.5), Inches(4.8), Inches(12.3), Inches(1.1))
add_text(s, "Training Transforms", Inches(0.65), Inches(4.88), Inches(3.5), Inches(0.38),
         size=16, bold=True, color=ACCENT)
add_text(s,
    "RandomCrop(256)  →  ToTensor()  →  Normalize( mean=[0.485, 0.456, 0.406], std=[0.229, 0.224, 0.225] )",
    Inches(0.65), Inches(5.3), Inches(12.0), Inches(0.38),
    size=14, color=BLUE_CODE)

# Bottom note
add_text(s, "SeededSampler enables exact mid-epoch resume — same image order guaranteed across interruptions",
         Inches(0.5), Inches(6.1), Inches(12.3), Inches(0.38),
         size=13, color=LIGHT_GRAY, italic=True)
slide_num(s, 11)


# ══════════════════════════════════════════════════════════════════════════════
# 12 — Style & Content Reference
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_bg(s); accent_bar(s)
title_text(s, "Style & Content Reference Images")

add_text(s, "Style Image  —  pencil_style_img1.jpg",
         Inches(0.5), Inches(1.1), Inches(5.5), Inches(0.45),
         size=17, bold=True, color=ACCENT)
add_image_safe(s, STYLE_IMG, Inches(0.5), Inches(1.65), Inches(5.5), Inches(4.8))

add_text(s, "Content Image  —  content_img1.png",
         Inches(6.3), Inches(1.1), Inches(5.5), Inches(0.45),
         size=17, bold=True, color=ACCENT)
add_image_safe(s, CONTENT_IMG, Inches(6.3), Inches(1.65), Inches(6.5), Inches(4.8))

add_rect(s, Inches(0.5), Inches(6.62), Inches(12.3), Inches(0.72), BG_CARD)
add_text(s,
    "⚙️  Full training in progress on MS-COCO train2014 (82,177 images, 2 epochs, batch 8)  ·  "
    "Stylized result images will replace this placeholder after training completes.",
    Inches(0.65), Inches(6.7), Inches(12.0), Inches(0.55),
    size=13, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)
slide_num(s, 12)


# ══════════════════════════════════════════════════════════════════════════════
# 13 — Performance: Classic vs Fast NST
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_bg(s); accent_bar(s)
title_text(s, "Classic NST vs Fast NST — Comparison")

headers = ["", "Classic NST", "Fast NST (ours)"]
rows = [
    ("Inference time",     "5–15 min",      "< 1 second"),
    ("Training required",  "None",          "~4–5 hrs (once)"),
    ("Quality",            "Very high",     "High (with improvements)"),
    ("New style",          "Re-run optim.", "Train new model"),
    ("GPU memory",         "~2 GB",         "~4 GB (training)"),
    ("Method",             "Pixel optim.",  "Feed-forward net"),
    ("Batch processing",   "❌ Impractical", "✅ Instant"),
    ("Adversarial loss",   "❌ No",          "✅ PatchGAN"),
    ("Best for",           "One-off art",   "Real-time / web"),
]
col_w = [Inches(3.8), Inches(4.0), Inches(4.0)]
col_x = [Inches(0.5), Inches(4.5), Inches(8.6)]

# Header row
for j, (hdr, cx, cw) in enumerate(zip(headers, col_x, col_w)):
    col = ACCENT if j == 2 else (ACCENT2 if j == 1 else BG_DARK)
    if j > 0:
        add_rect(s, cx, Inches(1.1), cw, Inches(0.55), col)
    add_text(s, hdr, cx + Inches(0.1), Inches(1.15), cw - Inches(0.2), Inches(0.42),
             size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER if j > 0 else PP_ALIGN.LEFT)

for i, row in enumerate(rows):
    y = Inches(1.72 + i * 0.6)
    bg = BG_CARD if i % 2 == 0 else RGBColor(0x12, 0x12, 0x25)
    for j, (val, cx, cw) in enumerate(zip(row, col_x, col_w)):
        add_rect(s, cx, y, cw, Inches(0.55), bg)
        col = WHITE if j == 0 else (GREEN_CODE if j == 2 else LIGHT_GRAY)
        bold = j == 0
        add_text(s, val, cx + Inches(0.12), y + Inches(0.09), cw - Inches(0.2), Inches(0.38),
                 size=13, bold=bold, color=col,
                 align=PP_ALIGN.CENTER if j > 0 else PP_ALIGN.LEFT)

slide_num(s, 13)


# ══════════════════════════════════════════════════════════════════════════════
# 14 — Future Work
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_bg(s); accent_bar(s)
title_text(s, "Future Work  — Currently in Active Development")

add_rect(s, Inches(0.5), Inches(1.0), Inches(12.3), Inches(0.48), RGBColor(0x1E,0x1A,0x0A))
add_text(s, "🚧  All features below are being actively built in this project",
         Inches(0.65), Inches(1.06), Inches(12.0), Inches(0.36),
         size=14, color=YELLOW, italic=True)

future = [
    ("⚡", "Fast NST Real-Time Inference",
     "TransformerNet trained on MS-COCO — sub-second stylization.\nReplaces 15-minute Classic NST for live use.", ACCENT),
    ("🎭", "PatchGAN Adversarial Training",
     "70×70 PatchDiscriminator with spectral norm.\nEnforces micro-texture realism — sharper brushstrokes.", ACCENT2),
    ("📐", "8 Training Improvements",
     "5-layer VGG loss · Multi-scale Gram · Feature stats loss\nIdentity loss · Per-layer weights · AMP + cosine LR", ACCENT3),
    ("🗂️", "Multi-Style Model Library",
     "train_all_styles.sh: one model per style image.\nInstant style switching from web UI dropdown.", ORANGE),
    ("📊", "TensorBoard Dashboard",
     "Live loss curves for all 5 loss components.\nSample preview images saved every 500 training steps.", PINK),
    ("🖼️", "Training Results in README",
     "Before/After images will be added to README.md\nafter the first full training run completes.", YELLOW),
]
for i, (icon, title, desc, col) in enumerate(future):
    c, row = i % 2, i // 2
    x = Inches(0.5 + c * 6.45)
    y = Inches(1.6 + row * 1.75)
    card(s, x, y, Inches(6.1), Inches(1.6), col)
    add_text(s, f"{icon}  {title}", x + Inches(0.22), y + Inches(0.1), Inches(5.65), Inches(0.42),
             size=15, bold=True, color=WHITE)
    add_text(s, desc, x + Inches(0.22), y + Inches(0.58), Inches(5.65), Inches(0.9),
             size=13, color=LIGHT_GRAY, wrap=True)
slide_num(s, 14)


# ══════════════════════════════════════════════════════════════════════════════
# 15 — Thank You
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_bg(s)
add_rect(s, 0, 0, Inches(0.12), SLIDE_H, ACCENT)

add_text(s, "Thank You",
         Inches(0.4), Inches(1.2), Inches(8.5), Inches(1.8),
         size=76, bold=True, color=WHITE)
add_text(s, "NST Studio  —  Neural Style Transfer Web Application",
         Inches(0.4), Inches(3.1), Inches(8.5), Inches(0.55),
         size=22, color=LIGHT_GRAY)

add_bullets(s, [
    "Classic NST: VGG-19 pixel optimization with live SSE streaming",
    "Full-stack Flask web app — drag-drop, gallery, comparison slider",
    "Fast NST in active development: < 1s inference per image",
    "PatchGAN adversarial training for texture realism",
    "Multi-style model library with one-command training scripts",
    "TensorBoard logging + sample previews every 500 training steps",
], Inches(0.4), Inches(3.85), Inches(8.5), Inches(3.2), size=16)

add_image_safe(s, STYLE_IMG, Inches(9.4), Inches(0.7), Inches(3.6), Inches(5.6))

add_text(s, "github.com/charbelmezeraani0-star  ·  2026",
         Inches(0.4), Inches(7.05), Inches(7), Inches(0.35),
         size=13, color=LIGHT_GRAY)
slide_num(s, 15)


# ── Save ──────────────────────────────────────────────────────────────────────
OUT = os.path.join(PROJECT_DIR, "NST_Studio_Presentation.pptx")
prs.save(OUT)
print(f"Saved: {OUT}")
